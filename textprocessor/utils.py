import PyPDF2
import os
import re
import logging
from typing import Dict, List, Optional, Tuple
import requests
from django.conf import settings
import json

logger = logging.getLogger(__name__)

def extract_text_from_file(file_path: str, file_type: str) -> Tuple[str, bool]:
    """
    Extract text from uploaded files
    Returns: (extracted_text, success)
    """
    try:
        if file_type == 'pdf':
            return extract_pdf_text(file_path)
        elif file_type == 'txt':
            return extract_txt_text(file_path)
        elif file_type == 'docx':
            return extract_docx_text(file_path)
        elif file_type == 'pptx':
            return extract_pptx_text(file_path)
        else:
            return '', False
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {str(e)}")
        return f"Error extracting text: {str(e)}", False

def extract_pdf_text(file_path: str) -> Tuple[str, bool]:
    """Extract text from PDF file"""
    try:
        text = ''
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text() or ''
                text += page_text + '\n'
        
        # Clean up the text
        text = clean_extracted_text(text)
        return text, True
    except Exception as e:
        logger.error(f"PDF extraction error: {str(e)}")
        return f"Error reading PDF: {str(e)}", False

def extract_txt_text(file_path: str) -> Tuple[str, bool]:
    """Extract text from TXT file"""
    try:
        # Try different encodings
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text = file.read()
                    return clean_extracted_text(text), True
            except UnicodeDecodeError:
                continue
        
        return "Error: Could not decode text file with any supported encoding", False
    except Exception as e:
        logger.error(f"TXT extraction error: {str(e)}")
        return f"Error reading text file: {str(e)}", False

def extract_docx_text(file_path: str) -> Tuple[str, bool]:
    """Extract text from DOCX file"""
    try:
        # Optional: requires python-docx
        # pip install python-docx
        try:
            from docx import Document
            doc = Document(file_path)
            text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            return clean_extracted_text(text), True
        except ImportError:
            logger.warning("python-docx not installed. Install with: pip install python-docx")
            return "Error: DOCX support not available. Please install python-docx", False
    except Exception as e:
        logger.error(f"DOCX extraction error: {str(e)}")
        return f"Error reading DOCX file: {str(e)}", False

def extract_pptx_text(file_path: str) -> Tuple[str, bool]:
    """Extract text from PPTX file"""
    try:
        # Optional: requires python-pptx
        # pip install python-pptx
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ''
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
            
            return clean_extracted_text(text), True
        except ImportError:
            logger.warning("python-pptx not installed. Install with: pip install python-pptx")
            return "Error: PPTX support not available. Please install python-pptx", False
    except Exception as e:
        logger.error(f"PPTX extraction error: {str(e)}")
        return f"Error reading PPTX file: {str(e)}", False

def clean_extracted_text(text: str) -> str:
    """Clean and normalize extracted text"""
    if not text:
        return ''
    
    # Remove extra whitespace and normalize line breaks
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n\s*\n', '\n\n', text)
    
    # Remove common PDF artifacts
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]', '', text)
    
    # Strip and return
    return text.strip()

def extract_topics_with_ai(text: str, max_topics: int = 20) -> Dict:
    """
    Use AI to extract topics and structure from educational content
    """
    if not text or len(text.strip()) < 100:
        return {"topics": [], "structure": {}, "summary": ""}
    
    # Truncate text if too long (AI API limits)
    if len(text) > 10000:
        text = text[:10000] + "..."
    
    prompt = f"""
    Analyze the following educational content and extract key topics, structure, and create a summary.
    
    Content:
    {text}
    
    Please respond with ONLY valid JSON in this exact format:
    {{
        "topics": ["topic1", "topic2", "topic3", ...],
        "structure": {{
            "chapters": ["chapter1", "chapter2", ...],
            "sections": {{
                "chapter1": ["section1", "section2", ...],
                "chapter2": ["section1", "section2", ...]
            }}
        }},
        "summary": "Brief summary of the content in 2-3 sentences",
        "suggested_categories": ["category1", "category2", ...]
    }}
    
    Extract up to {max_topics} main topics.
    """
    
    try:
        response = requests.post(
            "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent",
            headers={
                'Content-Type': 'application/json',
                'X-goog-api-key': getattr(settings, 'GEMINI_API_KEY', '')
            },
            json={
                "contents": [{"parts": [{"text": prompt}]}]
            },
            timeout=30
        )
        
        if response.status_code == 200:
            raw_text = (
                response.json()
                .get('candidates', [{}])[0]
                .get('content', {})
                .get('parts', [{}])[0]
                .get('text', '')
            )
            
            # Extract JSON from response
            json_match = re.search(r'(?:```json\s*)?(\{[\s\S]*?\})(?:\s*```)?', raw_text, re.DOTALL)
            if json_match:
                result = json.loads(json_match.group(1))
                return result
        
        logger.warning(f"AI topic extraction failed. Status: {response.status_code}")
        
    except Exception as e:
        logger.error(f"AI topic extraction error: {str(e)}")
    
    # Fallback: basic topic extraction
    return extract_topics_fallback(text, max_topics)

def extract_topics_fallback(text: str, max_topics: int = 20) -> Dict:
    """
    Fallback topic extraction using basic text processing
    """
    if not text:
        return {"topics": [], "structure": {}, "summary": ""}
    
    # Simple keyword extraction
    words = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', text)
    word_freq = {}
    
    for word in words:
        if len(word) > 3 and word.lower() not in ['the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can', 'had', 'her', 'was', 'one', 'our', 'out', 'day', 'get', 'has', 'him', 'his', 'how', 'man', 'new', 'now', 'old', 'see', 'two', 'way', 'who', 'boy', 'did', 'its', 'let', 'put', 'say', 'she', 'too', 'use']:
            word_freq[word] = word_freq.get(word, 0) + 1
    
    # Get top topics
    topics = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:max_topics]
    topic_list = [topic[0] for topic in topics]
    
    # Basic summary (first 2 sentences)
    sentences = re.split(r'[.!?]+', text)
    summary = '. '.join([s.strip() for s in sentences[:2] if s.strip()]) + '.'
    
    return {
        "topics": topic_list,
        "structure": {"chapters": [], "sections": {}},
        "summary": summary,
        "suggested_categories": []
    }

def get_text_chunk_for_quiz(text: str, topic: str = None, method: str = 'full_document') -> str:
    """
    Extract relevant text chunk for quiz generation
    """
    if not text:
        return ''
    
    if method == 'full_document':
        # Use entire document but limit to reasonable size
        return text[:8000] if len(text) > 8000 else text
    
    elif method == 'specific_topic' and topic:
        # Try to find and extract text related to specific topic
        topic_lower = topic.lower()
        text_lower = text.lower()
        
        # Find topic mentions
        topic_positions = []
        start = 0
        while True:
            pos = text_lower.find(topic_lower, start)
            if pos == -1:
                break
            topic_positions.append(pos)
            start = pos + 1
        
        if topic_positions:
            # Extract text around first mention (±2000 chars)
            pos = topic_positions[0]
            start = max(0, pos - 2000)
            end = min(len(text), pos + 2000)
            return text[start:end]
        else:
            # Topic not found, use full text
            return text[:8000] if len(text) > 8000 else text
    
    elif method == 'auto_summary':
        # Take chunks from beginning, middle, and end
        length = len(text)
        if length <= 6000:
            return text
        
        chunk_size = 2000
        beginning = text[:chunk_size]
        middle_start = (length - chunk_size) // 2
        middle = text[middle_start:middle_start + chunk_size]
        end = text[-chunk_size:]
        
        return f"{beginning}\n...\n{middle}\n...\n{end}"
    
    return text[:8000] if len(text) > 8000 else text

def validate_extracted_text(text: str) -> Dict[str, any]:
    """
    Validate and analyze extracted text quality
    """
    if not text:
        return {"valid": False, "reason": "No text extracted", "word_count": 0, "quality_score": 0}
    
    word_count = len(text.split())
    char_count = len(text)
    
    # Quality checks
    if word_count < 50:
        return {"valid": False, "reason": "Text too short (less than 50 words)", "word_count": word_count, "quality_score": 0}
    
    if char_count < 200:
        return {"valid": False, "reason": "Text too short (less than 200 characters)", "word_count": word_count, "quality_score": 0}
    
    # Calculate quality score (0-100)
    quality_score = min(100, (word_count / 100) * 50 + 50)  # Base score + word count bonus
    
    # Check for common extraction errors
    if text.count('?') > word_count * 0.1:  # Too many question marks might indicate OCR errors
        quality_score -= 20
    
    if re.search(r'[^\w\s.,!?;:()\-\'"]', text):  # Unusual characters
        quality_score -= 10
    
    quality_score = max(0, quality_score)
    
    return {
        "valid": True,
        "word_count": word_count,
        "char_count": char_count,
        "quality_score": quality_score,
        "reason": "Text extraction successful"
    }